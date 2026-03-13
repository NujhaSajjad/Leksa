import io
import logging

logger = logging.getLogger(__name__)


def parse_document(file_bytes: bytes, filename: str) -> str:
    """
    File extension dekh ke sahi parser call karo.
    Returns: plain text string
    """
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if ext == "pdf":
        return _parse_pdf(file_bytes)
    elif ext in ("ppt", "pptx"):
        return _parse_ppt(file_bytes)
    elif ext in ("doc", "docx"):
        return _parse_docx(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _parse_pdf(file_bytes: bytes) -> str:
    try:
        import fitz

        doc = fitz.open(stream=file_bytes, filetype="pdf")
        pages = []
        num_pages = len(doc)  # ← pehle yahan save karo

        for page_num, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                pages.append(f"[Page {page_num}]\n{text}")

        doc.close()
        result = "\n\n".join(pages)
        logger.info(f"PDF parsed: {num_pages} pages, {len(result)} chars")  # ← num_pages use karo
        return result


    except ImportError:
        raise RuntimeError("PyMuPDF install karo: pip install pymupdf")
    except Exception as e:
        raise RuntimeError(f"PDF parse error: {e}")


def _parse_ppt(file_bytes: bytes) -> str:
    """python-pptx se PPT/PPTX parse karo — slide by slide"""
    try:
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation(io.BytesIO(file_bytes))
        slides = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []

            # Slide title nikalo
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_texts.append(f"Title: {slide.shapes.title.text.strip()}")

            # Baaki sab shapes ka text
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            # Speaker notes bhi lo (teaching ke liye valuable)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"[Notes]: {notes}")

            if slide_texts:
                slides.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

        result = "\n\n".join(slides)
        logger.info(f"PPT parsed: {len(prs.slides)} slides, {len(result)} chars")
        return result

    except ImportError:
        raise RuntimeError("python-pptx install karo: pip install python-pptx")
    except Exception as e:
        raise RuntimeError(f"PPT parse error: {e}")


def _parse_docx(file_bytes: bytes) -> str:
    """python-docx se DOCX parse karo — paragraph by paragraph"""
    try:
        from docx import Document

        doc = Document(io.BytesIO(file_bytes))
        paragraphs = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Heading hai to mark karo
                if para.style.name.startswith("Heading"):
                    paragraphs.append(f"\n[{para.style.name}] {text}")
                else:
                    paragraphs.append(text)

        # Tables ka text bhi extract karo
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    paragraphs.append(row_text)

        result = "\n".join(paragraphs)
        logger.info(f"DOCX parsed: {len(doc.paragraphs)} paragraphs, {len(result)} chars")
        return result

    except ImportError:
        raise RuntimeError("python-docx install karo: pip install python-docx")
    except Exception as e:
        raise RuntimeError(f"DOCX parse error: {e}")
